"""PowerPoint (.pptx) export — premium, design-led deck (cover, section dividers,
2x2 SWOT matrix, cards, colour-coded priority badges, 5-year timeline, styled M&E table).
Flat modern style, institutional WHO palette."""
from __future__ import annotations
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from core.models import NISStrategy
from core.branding import logo_path

# ---- palette ----
DARK = RGBColor(0x00, 0x33, 0x66)
PRIMARY = RGBColor(0x00, 0x93, 0xD5)
ACCENT = RGBColor(0x6C, 0xB3, 0x3F)
BG = RGBColor(0xF4, 0xF7, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x20, 0x2A, 0x38)
MUTED = RGBColor(0x5B, 0x6B, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# SWOT quadrant colours (header / light fill)
S_C, S_BG = RGBColor(0x2E, 0x7D, 0x32), RGBColor(0xE9, 0xF5, 0xEA)
W_C, W_BG = RGBColor(0xC6, 0x28, 0x28), RGBColor(0xFC, 0xEA, 0xEA)
O_C, O_BG = RGBColor(0x15, 0x65, 0xC0), RGBColor(0xE7, 0xF0, 0xFB)
T_C, T_BG = RGBColor(0xEF, 0x6C, 0x00), RGBColor(0xFD, 0xF1, 0xE3)
PRIO = {"high": RGBColor(0x2E, 0x7D, 0x32), "medium": RGBColor(0xEF, 0x6C, 0x00),
        "low": RGBColor(0x90, 0x9C, 0xA8)}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def build_ppt(s: NISStrategy) -> bytes:
    fr = s.profile.language == "fr"
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    blank = prs.slide_layouts[6]
    period = f"{s.profile.nis_start_year}–{s.profile.nis_start_year + s.profile.nis_duration_years - 1}"
    years = s.profile.years
    lp = logo_path(s.profile.language)
    state = {"n": 0}

    # ---------- primitives ----------
    def slide(bg=BG):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb = bg
        return sl

    def rect(sl, x, y, w, h, fill, rounded=False, line=None, line_w=1.0):
        shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                                  x, y, w, h)
        shp.shadow.inherit = False
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        return shp

    def text(sl, s_, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False, spacing=1.04):
        tb = sl.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
        tf.word_wrap = True; tf.vertical_anchor = anchor
        items = s_ if isinstance(s_, list) else [s_]
        for i, line in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(2)
            r = p.add_run(); r.text = str(line)
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = color; r.font.name = "Calibri"
        return tb

    def bullets(sl, items, x, y, w, h, size=14, color=INK, bullet="•", maxn=8):
        tb = sl.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
        items = [i for i in items if i and str(i).strip()][:maxn] or ["—"]
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.12; p.space_after = Pt(5)
            r = p.add_run(); r.text = f"{bullet}  {it}"
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
        return tb

    def header(sl, kicker, title):
        rect(sl, 0, 0, Inches(0.18), Inches(1.25), PRIMARY)
        text(sl, kicker.upper(), Inches(0.55), Inches(0.28), Inches(11), Inches(0.3), 12, PRIMARY, bold=True)
        text(sl, title, Inches(0.55), Inches(0.55), Inches(12.2), Inches(0.8), 28, DARK, bold=True)
        footer(sl)

    def footer(sl):
        state["n"] += 1
        text(sl, f"SNV {s.profile.country_name} · {period}" if fr
             else f"NIS {s.profile.country_name} · {period}",
             Inches(0.55), Inches(7.02), Inches(8), Inches(0.3), 9, MUTED)
        text(sl, str(state["n"]), Inches(12.4), Inches(7.02), Inches(0.5), Inches(0.3), 9, MUTED,
             align=PP_ALIGN.RIGHT)

    def pill(sl, x, y, label, color):
        p = rect(sl, x, y, Inches(1.15), Inches(0.34), color, rounded=True)
        tf = p.text_frame; tf.word_wrap = False
        r = tf.paragraphs[0].add_run(); r.text = label
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ---------- slides ----------
    def title_slide():
        sl = slide(WHITE)
        rect(sl, 0, 0, EMU_W, Inches(0.32), PRIMARY)
        rect(sl, 0, Inches(6.0), EMU_W, Inches(1.5), DARK)
        if lp:
            try:
                sl.shapes.add_picture(str(lp), Inches(0.6), Inches(0.7), height=Inches(0.95))
            except Exception:
                pass
        text(sl, "Stratégie Nationale de Vaccination" if fr else "National Immunization Strategy",
             Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.6), 40, DARK, bold=True)
        rect(sl, Inches(0.85), Inches(3.95), Inches(2.2), Inches(0.08), ACCENT)
        text(sl, f"{s.profile.country_name}  ·  {period}", Inches(0.8), Inches(4.2),
             Inches(11), Inches(0.7), 22, PRIMARY, bold=True)
        text(sl, [s.profile.ministry_name, s.profile.epi_programme_name, s.profile.generation_date],
             Inches(0.8), Inches(6.15), Inches(11), Inches(1.2), 13, WHITE)

    def divider(num, title, sub=""):
        sl = slide(DARK)
        rect(sl, 0, Inches(3.2), EMU_W, Inches(0.08), ACCENT)
        text(sl, num, Inches(0.7), Inches(1.7), Inches(4), Inches(2), 110,
             RGBColor(0x1C, 0x4E, 0x82), bold=True)
        text(sl, title, Inches(0.8), Inches(3.45), Inches(11.6), Inches(1.2), 34, WHITE, bold=True)
        if sub:
            text(sl, sub, Inches(0.85), Inches(4.5), Inches(11.4), Inches(0.8), 16,
                 RGBColor(0xBF, 0xD6, 0xEC))

    def agenda():
        sl = slide()
        header(sl, "Sommaire" if fr else "Agenda", "Sommaire" if fr else "Agenda")
        items = [("Vision, but et objectif" if fr else "Vision, goal & objective"),
                 ("Analyse FFOM" if fr else "SWOT analysis"),
                 ("Causes profondes et obstacles" if fr else "Root causes & obstacles"),
                 ("Objectifs stratégiques" if fr else "Strategic objectives"),
                 ("Interventions prioritaires" if fr else "Priority interventions"),
                 ("Suivi et évaluation" if fr else "Monitoring & evaluation"),
                 ("Calendrier sur 5 ans" if fr else "5-year timeline"),
                 ("Conditions de succès & prochaines étapes" if fr else "Success factors & next steps")]
        cols, x0, y0, cw, ch, gx, gy = 2, Inches(0.55), Inches(1.7), Inches(6.0), Inches(1.0), Inches(0.3), Inches(0.18)
        for i, it in enumerate(items):
            cx = x0 + (i % cols) * (cw + gx); cy = y0 + (i // cols) * (ch + gy)
            rect(sl, cx, cy, cw, ch, CARD, rounded=True, line=RGBColor(0xE0, 0xE8, 0xF0))
            rect(sl, cx, cy, Inches(0.7), ch, PRIMARY, rounded=False)
            text(sl, str(i + 1), cx, cy, Inches(0.7), ch, 22, WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            text(sl, it, cx + Inches(0.85), cy, cw - Inches(1.0), ch, 15, INK, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    def cards(kicker, title, items):
        """items: list of (heading, body, accent_color)"""
        sl = slide()
        header(sl, kicker, title)
        n = max(1, len(items))
        x0, y0, gap = Inches(0.55), Inches(1.7), Inches(0.3)
        cw = (Inches(12.25) - gap * (n - 1)) / n
        for i, (hd, body, ac) in enumerate(items):
            cx = x0 + i * (cw + gap)
            rect(sl, cx, y0, cw, Inches(4.7), CARD, rounded=True, line=RGBColor(0xE0, 0xE8, 0xF0))
            rect(sl, cx, y0, cw, Inches(0.16), ac)
            text(sl, hd, cx + Inches(0.25), y0 + Inches(0.35), cw - Inches(0.5), Inches(0.7),
                 17, DARK, bold=True)
            text(sl, body or "—", cx + Inches(0.25), y0 + Inches(1.1), cw - Inches(0.5), Inches(3.4),
                 13, INK, spacing=1.12)

    def swot_matrix():
        sl = slide()
        header(sl, "Analyse FFOM" if fr else "SWOT analysis",
               "Forces, faiblesses, opportunités, menaces" if fr else
               "Strengths, weaknesses, opportunities, threats")
        agg = {"strengths": [], "weaknesses": [], "opportunities": [], "threats": []}
        for x in s.swot:
            for k in agg:
                for v in getattr(x, k):
                    if v and v not in agg[k]:
                        agg[k].append(v)
        quad = [("Forces" if fr else "Strengths", agg["strengths"], S_C, S_BG, 0, 0),
                ("Faiblesses" if fr else "Weaknesses", agg["weaknesses"], W_C, W_BG, 1, 0),
                ("Opportunités" if fr else "Opportunities", agg["opportunities"], O_C, O_BG, 0, 1),
                ("Menaces" if fr else "Threats", agg["threats"], T_C, T_BG, 1, 1)]
        x0, y0, cw, ch, g = Inches(0.55), Inches(1.65), Inches(6.05), Inches(2.55), Inches(0.2)
        for label, vals, col, bgc, cx, cy in quad:
            X = x0 + cx * (cw + g); Y = y0 + cy * (ch + g)
            rect(sl, X, Y, cw, ch, bgc, rounded=True, line=col, line_w=1.25)
            rect(sl, X, Y, Inches(0.14), ch, col)
            text(sl, label, X + Inches(0.3), Y + Inches(0.12), cw - Inches(0.6), Inches(0.4),
                 16, col, bold=True)
            bullets(sl, vals, X + Inches(0.3), Y + Inches(0.6), cw - Inches(0.55), ch - Inches(0.7),
                    size=11.5, color=INK, maxn=4)

    def interventions_slide():
        sl = slide()
        header(sl, "Interventions" if fr else "Interventions",
               "Interventions prioritaires" if fr else "Priority interventions")
        order = {"high": 0, "medium": 1, "low": 2}
        ivs = sorted(s.interventions,
                     key=lambda iv: order.get(getattr(iv.priority_level, "value", iv.priority_level), 3))[:5]
        y = Inches(1.7)
        for iv in ivs:
            lvl = getattr(iv.priority_level, "value", iv.priority_level)
            rect(sl, Inches(0.55), y, Inches(12.25), Inches(0.92), CARD, rounded=True,
                 line=RGBColor(0xE0, 0xE8, 0xF0))
            rect(sl, Inches(0.55), y, Inches(0.14), Inches(0.92), PRIO.get(lvl, MUTED))
            text(sl, iv.title, Inches(0.9), y + Inches(0.08), Inches(9.3), Inches(0.4), 15, DARK, bold=True)
            text(sl, (iv.expected_impact or iv.rationale or "")[:130], Inches(0.9), y + Inches(0.48),
                 Inches(9.3), Inches(0.4), 11, MUTED)
            pill(sl, Inches(11.4), y + Inches(0.28),
                 {"high": "ÉLEVÉ" if fr else "HIGH", "medium": "MOYEN" if fr else "MEDIUM",
                  "low": "FAIBLE" if fr else "LOW"}.get(lvl, str(lvl).upper()), PRIO.get(lvl, MUTED))
            y += Inches(1.04)
        if not ivs:
            text(sl, "—", Inches(0.6), Inches(2), Inches(6), Inches(0.5), 14, MUTED)

    def me_slide():
        sl = slide()
        header(sl, "Suivi & évaluation" if fr else "Monitoring & evaluation",
               "Cadre de suivi et d’évaluation" if fr else "M&E framework")
        inds = s.indicators[:6]
        cols = [("Indicateur" if fr else "Indicator", Inches(5.0)),
                ("Base" if fr else "Baseline", Inches(1.8))] + [(str(y_), Inches(0.95)) for y_ in years]
        rows = 1 + max(1, len(inds))
        tbl = sl.shapes.add_table(rows, len(cols), Inches(0.55), Inches(1.7),
                                  Inches(12.25), Inches(0.5) + Inches(0.7) * len(inds)).table
        for j, (h, w) in enumerate(cols):
            tbl.columns[j].width = w
            cell = tbl.cell(0, j); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = DARK
            r = cell.text_frame.paragraphs[0].runs[0]; r.font.size = Pt(10); r.font.bold = True
            r.font.color.rgb = WHITE
        for i, ind in enumerate(inds, 1):
            vals = [ind.name[:60], ind.baseline[:18]] + [ind.targets.get(f"Y{k+1}", "") for k in range(len(years))]
            for j, v in enumerate(vals):
                cell = tbl.cell(i, j); cell.text = str(v)
                pr = cell.text_frame.paragraphs[0].runs
                if pr:
                    pr[0].font.size = Pt(9); pr[0].font.color.rgb = INK
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else BG
        if not inds:
            text(sl, "—", Inches(0.6), Inches(2), Inches(5), Inches(0.5), 14, MUTED)

    def timeline_slide():
        sl = slide()
        header(sl, "Calendrier" if fr else "Timeline", "Calendrier stratégique sur 5 ans" if fr else
               "5-year strategic timeline")
        n = len(years)
        x0, y = Inches(0.9), Inches(4.0)
        span = Inches(11.6)
        rect(sl, x0, y, span, Inches(0.06), PRIMARY)
        step = span / max(1, n - 1) if n > 1 else span
        for i, yr in enumerate(years):
            cx = x0 + step * i
            dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.18), y - Inches(0.15),
                                      Inches(0.36), Inches(0.36))
            dot.fill.solid(); dot.fill.fore_color.rgb = DARK; dot.line.fill.background()
            cnt = sum(1 for a in s.activities if a.years.get(f"Y{i+1}"))
            up = (i % 2 == 0)
            text(sl, str(yr), cx - Inches(0.8), y + (Inches(0.35) if not up else Inches(-0.95)),
                 Inches(1.6), Inches(0.5), 18, DARK, bold=True, align=PP_ALIGN.CENTER)
            text(sl, f"{cnt} " + ("activités" if fr else "activities"),
                 cx - Inches(0.9), y + (Inches(0.85) if not up else Inches(-0.5)),
                 Inches(1.8), Inches(0.4), 11, MUTED, align=PP_ALIGN.CENTER)

    def closing():
        sl = slide(DARK)
        rect(sl, Inches(0.85), Inches(3.4), Inches(2.2), Inches(0.1), ACCENT)
        text(sl, "Merci" if fr else "Thank you", Inches(0.8), Inches(2.5), Inches(11.6), Inches(1),
             44, WHITE, bold=True)
        text(sl, ("Vers une vaccination équitable et durable" if fr
                  else "Towards equitable and sustainable immunization"),
             Inches(0.85), Inches(3.7), Inches(11), Inches(0.7), 18, RGBColor(0xBF, 0xD6, 0xEC))
        text(sl, f"{s.profile.ministry_name} · {s.profile.country_name} · {period}",
             Inches(0.85), Inches(6.6), Inches(11), Inches(0.5), 12, RGBColor(0x9F, 0xB8, 0xD2))

    # ---------- build ----------
    title_slide()
    agenda()
    divider("01", "Vision, but et objectif" if fr else "Vision, goal & objective")
    cards("Vision", "Vision, but et objectif" if fr else "Vision, goal & objective",
          [("Vision (≈10 ans)" if fr else "Vision (~10y)", s.vision.vision, PRIMARY),
           ("But de la SNV" if fr else "NIS goal", s.vision.goal, ACCENT),
           ("Objectif général" if fr else "Overall objective", s.vision.overall_objective, DARK)])
    divider("02", "Analyse FFOM" if fr else "SWOT analysis")
    swot_matrix()
    divider("03", "Causes profondes & objectifs" if fr else "Root causes & objectives")
    sl = slide(); header(sl, "Objectifs" if fr else "Objectives",
                         "Objectifs stratégiques prioritaires" if fr else "Priority strategic objectives")
    bullets(sl, [o.objective_text for o in s.objectives] or ["—"], Inches(0.6), Inches(1.7),
            Inches(12.2), Inches(5), size=15, bullet="◆", maxn=8)
    divider("04", "Interventions prioritaires" if fr else "Priority interventions")
    interventions_slide()
    divider("05", "Suivi & évaluation" if fr else "Monitoring & evaluation")
    me_slide()
    divider("06", "Calendrier & prochaines étapes" if fr else "Timeline & next steps")
    timeline_slide()
    cards("Conditions de succès" if fr else "Success factors",
          "Conditions de succès" if fr else "Conditions for success",
          [("Engagement & financement" if fr else "Commitment & financing",
            "Engagement politique et financement national durable." if fr
            else "Political commitment and sustainable national financing.", ACCENT),
           ("Coordination" if fr else "Coordination",
            "Partenaires alignés (OMS, UNICEF, Gavi) et coordination multisectorielle." if fr
            else "Aligned partners (WHO, UNICEF, Gavi) and multisectoral coordination.", PRIMARY),
           ("Données & redevabilité" if fr else "Data & accountability",
            "Données de qualité, suivi régulier et redevabilité à tous les niveaux." if fr
            else "Quality data, regular monitoring and accountability at all levels.", DARK)])
    sl = slide(); header(sl, "Prochaines étapes" if fr else "Next steps",
                         "Prochaines étapes" if fr else "Next steps")
    bullets(sl, ["Validation par le Ministère de la Santé" if fr else "MoH validation",
                 "Chiffrage de la stratégie (NIS.COST)" if fr else "Strategy costing (NIS.COST)",
                 "Plan opérationnel annuel (POA)" if fr else "Annual operational plan (AOP)",
                 "Mobilisation des ressources" if fr else "Resource mobilization"],
            Inches(0.6), Inches(1.8), Inches(12), Inches(4.5), size=18, bullet="→", maxn=6)
    closing()

    buf = BytesIO(); prs.save(buf)
    return buf.getvalue()
