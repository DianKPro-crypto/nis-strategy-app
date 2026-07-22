"""
Document ingestion: extract text, tables, metadata and locators (page/slide/sheet)
from .docx / .xlsx / .pptx / .pdf / .csv / .txt files.

Each extractor is defensive: a corrupted or unreadable file yields an
UploadedDocument with an error note rather than crashing the app.
"""
from __future__ import annotations
import csv
import io
import re
from pathlib import Path

from config.settings import ALLOWED_EXTENSIONS, MAX_TEXT_CHARS_PER_DOC
from core.models import UploadedDocument


def extract_document(file_bytes: bytes, filename: str, doc_category: str = "") -> UploadedDocument:
    ext = Path(filename).suffix.lower().lstrip(".")
    doc = UploadedDocument(name=filename, file_type=ext, size_bytes=len(file_bytes),
                           doc_category=doc_category)
    if ext not in ALLOWED_EXTENSIONS:
        doc.text = f"[Type de fichier non pris en charge: .{ext}]"
        return doc
    try:
        if ext == "txt":
            _txt(file_bytes, doc)
        elif ext == "csv":
            _csv(file_bytes, doc)
        elif ext == "docx":
            _docx(file_bytes, doc)
        elif ext == "xlsx":
            _xlsx(file_bytes, doc)
        elif ext == "pptx":
            _pptx(file_bytes, doc)
        elif ext == "pdf":
            _pdf(file_bytes, doc)
    except Exception as e:  # never let one bad file break ingestion
        doc.text = (doc.text or "") + f"\n[Erreur d’extraction: {e}]"
        doc.metadata["extraction_error"] = str(e)
    doc.text = (doc.text or "")[:MAX_TEXT_CHARS_PER_DOC]
    return doc


def _txt(b: bytes, doc: UploadedDocument) -> None:
    doc.text = b.decode("utf-8", errors="replace")


def _csv(b: bytes, doc: UploadedDocument) -> None:
    text = b.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    doc.text = "\n".join("\t".join(r) for r in rows[:500])
    doc.tables_summary = f"[CSV: {len(rows)} lignes, {len(rows[0]) if rows else 0} colonnes]"
    doc.metadata["rows"] = len(rows)


def _docx(b: bytes, doc: UploadedDocument) -> None:
    import docx  # python-docx
    d = docx.Document(io.BytesIO(b))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    n_tables = len(d.tables)
    for ti, table in enumerate(d.tables, 1):
        parts.append(f"[Tableau {ti}]")
        for row in table.rows:
            parts.append(" | ".join(c.text for c in row.cells))
    doc.text = "\n".join(parts)
    doc.tables_summary = f"[Word: {len(d.paragraphs)} paragraphes, {n_tables} tableaux]"
    doc.metadata["tables"] = n_tables


def _xlsx(b: bytes, doc: UploadedDocument) -> None:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(b), data_only=True, read_only=True)
    parts, sheets = [], []
    for ws in wb.worksheets:
        sheets.append(ws.title)
        parts.append(f"=== FEUILLE: {ws.title} ===")
        rows = list(ws.iter_rows(values_only=True))
        # Detect a header row = the row (in the first 8) with the most distinct text cells.
        # This lets us LABEL each value with its column (e.g. "POURQUOI-1", "Problème principal"),
        # so the structure of a WHO workbook survives extraction instead of being flattened.
        def _score(r):
            return sum(1 for c in r if isinstance(c, str) and c.strip())
        header, h_idx = None, -1
        for i, r in enumerate(rows[:8]):
            if _score(r) >= 3 and _score(r) > (_score(header) if header else 0):
                header, h_idx = r, i
        hdr = [str(c).strip() if c is not None else "" for c in header] if header else []
        for i, row in enumerate(rows):
            if i > 300:
                parts.append("... (tronqué)")
                break
            vals = [("" if c is None else str(c).strip()) for c in row]
            if not any(vals):
                continue
            if i == h_idx:
                parts.append("COLONNES: " + " | ".join(v for v in hdr if v))
                continue
            if hdr:
                # Pair each non-empty value with its column header -> "Colonne: valeur".
                pairs = []
                for ci, v in enumerate(vals):
                    if not v:
                        continue
                    col = hdr[ci] if ci < len(hdr) and hdr[ci] else f"col{ci+1}"
                    pairs.append(f"{col}: {v[:600]}")
                parts.append(" | ".join(pairs))
            else:
                parts.append("\t".join(v for v in vals if v))
    doc.text = "\n".join(parts)
    doc.tables_summary = f"[Excel: feuilles = {', '.join(sheets)}]"
    doc.metadata["sheets"] = sheets
    # Structured parse of a completed WHO "sequence of events" sheet so columns K/L
    # (Dernier POURQUOI / Problème principal) are reused VERBATIM, not paraphrased by the AI.
    try:
        seq = _parse_who_sequence(wb)
        if seq:
            doc.metadata["who_sequence"] = seq
    except Exception:
        pass


_SUBCODE_RE = re.compile(r"^\s*(\d+\.\d+)\b")


def _parse_who_sequence(wb) -> dict:
    """Map each sub-component code -> its country-filled analysis from the WHO 'Séquence' sheet:
    {'1.1': {'probleme_principal': str, 'dernier_pourquoi': str, 'weaknesses': [str,...]}}."""
    def _norm(x):
        return (str(x).strip() if x is not None else "")

    for ws in wb.worksheets:
        rows = [[ _norm(c) for c in r] for r in ws.iter_rows(values_only=True)]
        # locate the header row that carries the WHO columns
        h_idx = -1
        for i, r in enumerate(rows[:8]):
            joined = " | ".join(r)
            if "POURQUOI" in joined and ("Problème principal" in joined or "Problème/Obstacle" in joined):
                h_idx = i; break
        if h_idx < 0:
            continue
        hdr = rows[h_idx]

        def _find(*needles):
            for ci, h in enumerate(hdr):
                if any(n.lower() in h.lower() for n in needles):
                    return ci
            return None
        c_sub = _find("composantes et sous", "sous-composante") or 0
        c_last = _find("Dernier POURQUOI")
        c_prob = _find("Problème principal", "Problème/Obstacle")
        c_weak = _find("Faiblesse")
        if c_prob is None and c_last is None:
            continue

        out: dict = {}
        cur = None
        for r in rows[h_idx + 1:]:
            cell0 = r[c_sub] if c_sub < len(r) else ""
            m = _SUBCODE_RE.match(cell0)
            if m:
                cur = m.group(1)
                out.setdefault(cur, {"probleme_principal": "", "dernier_pourquoi": [], "weaknesses": []})
            if cur is None:
                continue
            def _cell(idx):
                return r[idx] if (idx is not None and idx < len(r)) else ""
            prob = _cell(c_prob)
            if prob and prob not in ("0",) and not out[cur]["probleme_principal"]:
                out[cur]["probleme_principal"] = prob
            last = _cell(c_last)
            if last and last not in ("0",):
                out[cur]["dernier_pourquoi"].append(last)
            weak = _cell(c_weak)
            if weak and weak not in ("0",):
                out[cur]["weaknesses"].append(weak)
        # keep only sub-components that actually carry country content
        clean = {}
        for code, v in out.items():
            if v["probleme_principal"] or v["dernier_pourquoi"]:
                # Bullet each distinct obstacle on its own line for readability (col K).
                obstacles = [o for o in v["dernier_pourquoi"] if o]
                dernier = ("• " + "\n• ".join(obstacles)) if len(obstacles) > 1 else "".join(obstacles)
                clean[code] = {"probleme_principal": v["probleme_principal"],
                               "dernier_pourquoi": dernier,
                               "weaknesses": v["weaknesses"]}
        if clean:
            return clean
    return {}


def _pptx(b: bytes, doc: UploadedDocument) -> None:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(b))
    parts = []
    for si, slide in enumerate(prs.slides, 1):
        parts.append(f"=== DIAPOSITIVE {si} ===")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
    doc.n_pages = len(prs.slides)
    doc.text = "\n".join(parts)
    doc.tables_summary = f"[PowerPoint: {doc.n_pages} diapositives]"


def _pdf(b: bytes, doc: UploadedDocument) -> None:
    text = ""
    # Prefer PyMuPDF (fitz); fall back to pdfplumber.
    try:
        import fitz  # PyMuPDF
        pdf = fitz.open(stream=b, filetype="pdf")
        pages = [f"[Page {i+1}]\n{pg.get_text()}" for i, pg in enumerate(pdf)]
        doc.n_pages = pdf.page_count
        text = "\n".join(pages)
    except Exception:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(b)) as pdf:
            pages = [f"[Page {i+1}]\n{(pg.extract_text() or '')}" for i, pg in enumerate(pdf.pages)]
            doc.n_pages = len(pdf.pages)
            text = "\n".join(pages)
    doc.text = text
    doc.tables_summary = f"[PDF: {doc.n_pages} pages]"
